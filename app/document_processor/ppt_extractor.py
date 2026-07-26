from pptx import Presentation

def extract_text_from_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_text.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(cells))
        text_parts.append(f"[Diapositiva {slide_num}]\n" + "\n".join(slide_text))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_parts.append(f"[Notas diapositiva {slide_num}]: {notes}")
    return "\n\n".join(text_parts)
